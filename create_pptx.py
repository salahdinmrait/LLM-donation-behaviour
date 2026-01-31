from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(15, 23, 42)       # #0f172a
PRIMARY = RGBColor(37, 99, 235)       # #2563eb
ACCENT = RGBColor(6, 182, 212)        # #06b6d4
SECONDARY = RGBColor(124, 58, 237)    # #7c3aed
LIGHT = RGBColor(248, 250, 252)       # #f8fafc
GRAY = RGBColor(100, 116, 139)        # #64748b
GREEN = RGBColor(34, 197, 94)         # #22c55e
RED = RGBColor(239, 68, 68)           # #ef4444

def add_dark_slide(prs):
    """Add a blank slide with dark background"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add dark background shape
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    return slide

def add_title_text(slide, text, top, font_size=44, color=LIGHT, bold=True, width=None):
    """Add a title text box"""
    left = Inches(0.8)
    if width is None:
        width = Inches(11.7)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    return txBox

def add_body_text(slide, text, left, top, width, height, font_size=18, color=LIGHT, bold=False):
    """Add body text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    return txBox

def add_bullet_point(slide, left, top, number, text, width=Inches(5.5)):
    """Add a numbered bullet point"""
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    
    # Number text
    tf = circle.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = LIGHT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bullet text
    add_body_text(slide, text, left + Inches(0.55), top, width, Inches(0.8), font_size=18)

def add_stat_box(slide, left, top, number, label, width=Inches(3)):
    """Add a statistics box"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(37, 99, 235)
    box.fill.fore_color.brightness = 0.8
    box.line.color.rgb = PRIMARY
    
    # Number
    num_box = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.7))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    label_box = slide.shapes.add_textbox(left, top + Inches(0.9), width, Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_image_placeholder(slide, left, top, width, height, image_path):
    """Add an image if it exists, otherwise add a placeholder"""
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width, height)
    else:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.color.rgb = PRIMARY

# ============== SLIDE 1: Title ==============
slide = add_dark_slide(prs)

# Main title
title = add_title_text(slide, 
    "Gender Differences in Responses to\nLLM-Generated vs Human-Written\nDonation Appeals",
    Inches(1.5), font_size=44)
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Subtitle
sub = add_body_text(slide, "A Multi-factorial Analysis of Digital Charitable Giving",
    Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.5), font_size=24, color=ACCENT)
sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Author info
info = add_body_text(slide, 
    "Salah-din Mrait\nBachelor End Project • Data Science Joint Degree\nTilburg University — Eindhoven University of Technology\n\nSupervisor: John Caffier\nJanuary 2026",
    Inches(0.8), Inches(4.3), Inches(11.7), Inches(2), font_size=18, color=GRAY)
info.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============== SLIDE 2: Research Problem ==============
slide = add_dark_slide(prs)
add_title_text(slide, "The Research Problem", Inches(0.5), font_size=40, color=ACCENT)

add_bullet_point(slide, Inches(0.8), Inches(1.5), "1", 
    "Charitable organizations increasingly use Large Language Models (LLMs) to create donation appeals")
add_bullet_point(slide, Inches(0.8), Inches(2.3), "2", 
    "Women donate more frequently and give higher amounts than men across cultures")
add_bullet_point(slide, Inches(0.8), Inches(3.1), "3", 
    "LLMs produce emotionally polished language but lack genuine emotional experience")
add_bullet_point(slide, Inches(0.8), Inches(3.9), "4", 
    "Will women — who value emotional authenticity — respond differently to LLM content?")

# Research question box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(5.5), Inches(3))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(30, 41, 59)
box.line.color.rgb = PRIMARY

add_body_text(slide, "Main Research Question", Inches(7.3), Inches(2), Inches(5), Inches(0.5), font_size=18, color=ACCENT, bold=True)
add_body_text(slide, '"Do women and men respond differently to LLM-generated versus human-written emotional donation appeals for cancer charities?"',
    Inches(7.3), Inches(2.5), Inches(5), Inches(2), font_size=16, color=LIGHT)

# ============== SLIDE 3: Theoretical Background ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Theoretical Background", Inches(0.5), font_size=40, color=ACCENT)

# Theory boxes
y_pos = Inches(1.5)
for emoji, title, desc in [
    ("🧠", "Empathic Concern", "Women report higher empathy. Perceived sincerity matters when emotional engagement drives persuasion."),
    ("⚠️", "Algorithm Aversion", "People trust algorithms less in emotional domains. Women may react more negatively to disclosed LLM content."),
    ("❤️", "Care vs. Justice", "Women adopt care-focused moral reasoning; LLM appeals may fail if caring relationship perception weakens.")
]:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.7), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    box.line.fill.background()
    
    add_body_text(slide, emoji, Inches(1), y_pos + Inches(0.4), Inches(0.5), Inches(0.5), font_size=28)
    add_body_text(slide, title, Inches(1.6), y_pos + Inches(0.2), Inches(3), Inches(0.4), font_size=20, color=ACCENT, bold=True)
    add_body_text(slide, desc, Inches(1.6), y_pos + Inches(0.6), Inches(10), Inches(0.6), font_size=16, color=LIGHT)
    
    y_pos += Inches(1.5)

# ============== SLIDE 4: Hypotheses ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Research Hypotheses", Inches(0.5), font_size=40, color=ACCENT)

hypotheses = [
    ("H1", "Women will report significantly lower persuasiveness scores (1-7 Likert) for LLM appeals, while men show no difference."),
    ("H2a", "Women will allocate higher monetary amounts (USD 0.00-0.10) regardless of appeal source."),
    ("H2b", "The donation gap between genders will be larger for human-written appeals."),
    ("H3", "Women's likelihood of selecting 'Like' will be higher for human-written appeals.")
]

y_pos = Inches(1.3)
for label, text in hypotheses:
    # Left border accent
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_pos, Inches(0.1), Inches(1.1))
    border.fill.solid()
    border.fill.fore_color.rgb = SECONDARY
    border.line.fill.background()
    
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), y_pos, Inches(11.6), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    box.line.fill.background()
    
    add_body_text(slide, label, Inches(1.1), y_pos + Inches(0.15), Inches(1), Inches(0.3), font_size=16, color=SECONDARY, bold=True)
    add_body_text(slide, text, Inches(1.1), y_pos + Inches(0.45), Inches(11), Inches(0.6), font_size=16, color=LIGHT)
    
    y_pos += Inches(1.3)

# ============== SLIDE 5: Methods ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Research Methods", Inches(0.5), font_size=40, color=ACCENT)

# Stats boxes
add_stat_box(slide, Inches(1.5), Inches(1.5), "725", "U.S. Participants")
add_stat_box(slide, Inches(5.2), Inches(1.5), "4,350", "Total Observations")
add_stat_box(slide, Inches(8.9), Inches(1.5), "6", "Appeals per Participant")

# Method details
add_body_text(slide, "📊 Design: Within-subjects experiment • Human vs LLM-generated appeals • Cancer charities",
    Inches(0.8), Inches(3.8), Inches(12), Inches(0.5), font_size=20, color=LIGHT)
add_body_text(slide, "📏 Outcomes: Engagement (Dislike/Neutral/Like) • Persuasiveness (1-7 scale) • Donation ($0.00-$0.10)",
    Inches(0.8), Inches(4.4), Inches(12), Inches(0.5), font_size=20, color=LIGHT)

# ============== SLIDE 6: Analytical Approach ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Analytical Approach", Inches(0.5), font_size=40, color=ACCENT)

add_bullet_point(slide, Inches(0.8), Inches(1.5), "1", 
    "Engagement: Ordinal logistic mixed-effects regression (cumulative link model)")
add_bullet_point(slide, Inches(0.8), Inches(2.3), "2", 
    "Persuasiveness: Linear mixed-effects model (continuous outcome)")
add_bullet_point(slide, Inches(0.8), Inches(3.1), "3", 
    "Donation: Hurdle model (35.3% zero-inflation required two-part approach)")

# Model features box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.5), Inches(5), Inches(2.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(30, 41, 59)
box.line.color.rgb = PRIMARY

add_body_text(slide, "All models include:", Inches(7.8), Inches(1.7), Inches(4.5), Inches(0.4), font_size=18, color=ACCENT)
add_body_text(slide, "✓ Random intercepts for participants\n✓ Random intercepts for posts\n✓ Gender × Content Source interaction",
    Inches(7.8), Inches(2.2), Inches(4.5), Inches(1.5), font_size=16, color=LIGHT)
# ============== SLIDE 7: Model Justification ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Model Selection Justification", Inches(0.5), font_size=40, color=ACCENT)

# Q-Q Plot for Persuasiveness
base_path = "thesis_source/thesis_latex/phase1_plots/"
add_image_placeholder(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.8),
    base_path + "persuasiveness_qqplot.png")
add_body_text(slide, "Persuasiveness: Q-Q plot confirms normality → Linear mixed-effects model",
    Inches(0.8), Inches(5.5), Inches(5.5), Inches(0.5), font_size=14, color=GRAY)

# Donation Histogram
add_image_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(3.8),
    base_path + "donation_histogram.png")
add_body_text(slide, "Donation: 35.3% zeros → Hurdle model required",
    Inches(6.8), Inches(5.5), Inches(5.5), Inches(0.5), font_size=14, color=GRAY)

# ============== SLIDE 8: Descriptive Results ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Descriptive Statistics", Inches(0.5), font_size=40, color=ACCENT)

# Image placeholders
base_path = "thesis_source/thesis_latex/phase1_plots/"
add_image_placeholder(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(4.5), 
    base_path + "engagement_distribution.png")
add_image_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.5), 
    base_path + "persuasiveness_boxplot.png")

# ============== SLIDE 9: Main Effects ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Main Effects: The LLM Advantage", Inches(0.5), font_size=40, color=ACCENT)

# Table header
headers = ["Outcome", "LLM Effect", "p-value", "Interpretation"]
x_positions = [Inches(0.8), Inches(3.5), Inches(5.5), Inches(7.5)]
widths = [Inches(2.5), Inches(2), Inches(2), Inches(5)]

y = Inches(1.5)
for i, header in enumerate(headers):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_positions[i], y, widths[i], Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(37, 99, 235)
    box.line.fill.background()
    add_body_text(slide, header, x_positions[i] + Inches(0.1), y + Inches(0.1), widths[i], Inches(0.4), font_size=16, color=LIGHT, bold=True)

# Table rows
rows = [
    ("Engagement", "OR = 1.58", "< .001", "58% higher odds of favorable rating"),
    ("Persuasiveness", "b = 0.20", "< .001", "+0.20 points on 7-point scale"),
    ("Donation Amount", "b = 0.0018", ".011", "Higher amounts for LLM appeals")
]

y = Inches(2.0)
for row in rows:
    for i, cell in enumerate(row):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_positions[i], y, widths[i], Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.color.rgb = RGBColor(50, 60, 80)
        color = GREEN if i == 2 else LIGHT
        add_body_text(slide, cell, x_positions[i] + Inches(0.1), y + Inches(0.15), widths[i], Inches(0.4), font_size=14, color=color)
    y += Inches(0.6)

# Key finding
add_body_text(slide, "LLM-generated appeals outperformed human-written appeals on ALL outcomes",
    Inches(0.8), Inches(4.5), Inches(12), Inches(0.5), font_size=24, color=ACCENT, bold=True)

# ============== SLIDE 10: Gender Effects on Donation ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Gender Effects on Donation", Inches(0.5), font_size=40, color=ACCENT)

# Donation means graph
add_image_placeholder(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.2),
    base_path + "donation_means.png")

# Stat box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.3), Inches(1.3))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(30, 41, 59)
box.line.color.rgb = PRIMARY

# Number
num_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.3), Inches(0.7))
tf = num_box.text_frame
p = tf.paragraphs[0]
p.text = "$0.0013"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

# Label
label_box = slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.3), Inches(0.4))
tf = label_box.text_frame
p = tf.paragraphs[0]
p.text = "Women donate more (p = .026)"
p.font.size = Pt(14)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# Points
add_body_text(slide, "✓ H2a Supported: Women allocate higher amounts than men",
    Inches(7), Inches(3.2), Inches(5.5), Inches(0.5), font_size=16, color=ACCENT)
add_body_text(slide, "✗ Gender did NOT predict whether someone donated at all (p = .412)",
    Inches(7), Inches(3.8), Inches(5.5), Inches(0.5), font_size=16, color=LIGHT)

# ============== SLIDE 11: Engagement Interaction ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Interaction: Engagement", Inches(0.5), font_size=40, color=ACCENT)
add_body_text(slide, "Gender × Content Source (OR = 1.53, p = .011) — SIGNIFICANT",
    Inches(0.8), Inches(1.0), Inches(8), Inches(0.4), font_size=18, color=GREEN)

# Graph
add_image_placeholder(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.2),
    "thesis_source/thesis_latex/phase3_plots/interaction_engagement_like.png")

# Interpretation
add_body_text(slide, "→ Both genders prefer LLM content",
    Inches(6.8), Inches(2), Inches(5.5), Inches(0.4), font_size=18, color=LIGHT)
add_body_text(slide, "→ Women show stronger preference for LLM appeals",
    Inches(6.8), Inches(2.6), Inches(5.5), Inches(0.4), font_size=18, color=ACCENT)
add_body_text(slide, "→ Steeper slope for women = significant interaction",
    Inches(6.8), Inches(3.2), Inches(5.5), Inches(0.4), font_size=18, color=LIGHT)

# Result box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4), Inches(5.5), Inches(1))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(60, 30, 30)
box.line.fill.background()
add_body_text(slide, "H3 NOT Supported: Women preferred LLM content more, not less",
    Inches(7), Inches(4.3), Inches(5), Inches(0.6), font_size=16, color=RED)

# ============== SLIDE 10: Persuasiveness Interaction ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Interaction: Persuasiveness", Inches(0.5), font_size=40, color=ACCENT)
add_body_text(slide, "Gender × Content Source (b = 0.06, p = .428) — NOT SIGNIFICANT",
    Inches(0.8), Inches(1.0), Inches(8), Inches(0.4), font_size=18, color=RED)

add_image_placeholder(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.2),
    "thesis_source/thesis_latex/phase3_plots/interaction_persuasiveness.png")

add_body_text(slide, "→ Both rated LLM content as more persuasive",
    Inches(6.8), Inches(2), Inches(5.5), Inches(0.4), font_size=18, color=LIGHT)
add_body_text(slide, "→ Visible gap is not statistically significant",
    Inches(6.8), Inches(2.6), Inches(5.5), Inches(0.4), font_size=18, color=LIGHT)

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4), Inches(5.5), Inches(1))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(60, 30, 30)
box.line.fill.background()
add_body_text(slide, "H1 NOT Supported: No interaction on persuasiveness",
    Inches(7), Inches(4.3), Inches(5), Inches(0.6), font_size=16, color=RED)

# ============== SLIDE 11: Donation Interaction ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Interaction: Donation Amount", Inches(0.5), font_size=40, color=ACCENT)
add_body_text(slide, "Gender × Content Source (b = 0.0002, p = .735) — NOT SIGNIFICANT",
    Inches(0.8), Inches(1.0), Inches(8), Inches(0.4), font_size=18, color=RED)

add_image_placeholder(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.2),
    "thesis_source/thesis_latex/phase3_plots/interaction_donation_fixed.png")

add_body_text(slide, "→ Women donate more across both conditions",
    Inches(6.8), Inches(2), Inches(5.5), Inches(0.4), font_size=18, color=LIGHT)
add_body_text(slide, "→ Parallel lines = no interaction",
    Inches(6.8), Inches(2.6), Inches(5.5), Inches(0.4), font_size=18, color=LIGHT)

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4), Inches(5.5), Inches(1))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(60, 30, 30)
box.line.fill.background()
add_body_text(slide, "H2b NOT Supported: Gender gap consistent across sources",
    Inches(7), Inches(4.3), Inches(5), Inches(0.6), font_size=16, color=RED)

# ============== SLIDE 12: Hypothesis Summary ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Hypothesis Testing Summary", Inches(0.5), font_size=40, color=ACCENT)

# Table
headers = ["Hypothesis", "Prediction", "Result"]
x_positions = [Inches(0.8), Inches(2.5), Inches(8)]
widths = [Inches(1.5), Inches(5.5), Inches(4)]

y = Inches(1.5)
for i, header in enumerate(headers):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_positions[i], y, widths[i], Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY
    box.line.fill.background()
    add_body_text(slide, header, x_positions[i] + Inches(0.1), y + Inches(0.1), widths[i], Inches(0.4), font_size=16, color=LIGHT, bold=True)

rows = [
    ("H1", "Women rate LLM as less persuasive", "Not Supported", RED),
    ("H2a", "Women donate more than men", "Supported", GREEN),
    ("H2b", "Larger gender gap for human appeals", "Not Supported", RED),
    ("H3", "Women prefer human content for engagement", "Not Supported (Opposite)", RED)
]

y = Inches(2.0)
for h, pred, result, result_color in rows:
    for i, (cell, color) in enumerate([(h, LIGHT), (pred, LIGHT), (result, result_color)]):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_positions[i], y, widths[i], Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.color.rgb = RGBColor(50, 60, 80)
        add_body_text(slide, cell, x_positions[i] + Inches(0.1), y + Inches(0.2), widths[i] - Inches(0.2), Inches(0.4), font_size=14, color=color)
    y += Inches(0.7)

# ============== SLIDE 13: Discussion ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Discussion", Inches(0.5), font_size=40, color=ACCENT)

add_bullet_point(slide, Inches(0.8), Inches(1.5), "1", 
    '"Sincerity Gap": Gender differences appeared only for quick, automatic engagement — not for deeper judgments or actual behavior.')
add_bullet_point(slide, Inches(0.8), Inches(2.5), "2", 
    "Algorithm Aversion Not Activated: Participants weren't told about LLM authorship, so skepticism wasn't triggered.")
add_bullet_point(slide, Inches(0.8), Inches(3.5), "3", 
    "Care Orientation Persists: Women's higher donations are triggered by distress signals — regardless of human or machine authorship.")

# ============== SLIDE 14: Practical Implications ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Practical Implications", Inches(0.5), font_size=40, color=ACCENT)

add_body_text(slide, "✓ LLMs can generate appeals that perform at least as well as human-written content",
    Inches(0.8), Inches(1.5), Inches(7), Inches(0.5), font_size=20, color=LIGHT)
add_body_text(slide, "✓ Especially useful for nonprofits with limited resources",
    Inches(0.8), Inches(2.1), Inches(7), Inches(0.5), font_size=20, color=LIGHT)
add_body_text(slide, "⚠️ Disclosure policies require caution — labeling may reduce impact",
    Inches(0.8), Inches(2.7), Inches(7), Inches(0.5), font_size=20, color=LIGHT)
add_body_text(slide, "→ Organizations face a trade-off between transparency and effectiveness",
    Inches(0.8), Inches(3.3), Inches(7), Inches(0.5), font_size=20, color=LIGHT)

# Key takeaway box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.5), Inches(5), Inches(3))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(30, 50, 70)
box.line.color.rgb = PRIMARY

add_body_text(slide, "Key Takeaway:", Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.4), font_size=18, color=LIGHT, bold=True)
add_body_text(slide, "Gender does not limit the value of LLMs in fundraising. Donors respond to content quality, not author identity.",
    Inches(7.8), Inches(2.3), Inches(4.5), Inches(2), font_size=18, color=ACCENT)

# ============== SLIDE 15: Limitations ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Limitations & Future Research", Inches(0.5), font_size=40, color=ACCENT)

limitations = [
    "🔒 Participants not disclosed about LLM authorship",
    "🎯 Single domain (cancer charities)",
    "💰 Small donation amounts ($0.00-$0.10)",
    "🇺🇸 U.S. sample only"
]

y = Inches(1.5)
for lim in limitations:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(5.5), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    box.line.fill.background()
    add_body_text(slide, lim, Inches(1), y + Inches(0.15), Inches(5), Inches(0.4), font_size=16, color=LIGHT)
    y += Inches(0.7)

# Future directions
add_body_text(slide, "Future Directions", Inches(7), Inches(1.5), Inches(5), Inches(0.4), font_size=20, color=ACCENT, bold=True)
add_body_text(slide, "1. Test effects when LLM authorship is disclosed\n2. Explore other charitable domains\n3. Examine long-term donor relationships",
    Inches(7), Inches(2), Inches(5.5), Inches(2), font_size=16, color=LIGHT)

# ============== SLIDE 16: Conclusion ==============
slide = add_dark_slide(prs)
add_title_text(slide, "Conclusion", Inches(0.5), font_size=40, color=ACCENT)

# Main conclusion box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(30, 50, 70)
box.line.color.rgb = PRIMARY

add_body_text(slide, "LLM-generated donation appeals outperform human-written appeals across all outcomes.",
    Inches(2), Inches(2), Inches(9), Inches(0.8), font_size=22, color=LIGHT)

add_body_text(slide, "Gender differences do not limit the effectiveness of LLM-generated charitable appeals.",
    Inches(2), Inches(3), Inches(9), Inches(0.8), font_size=26, color=ACCENT, bold=True)

add_body_text(slide, "Women's higher generosity persists regardless of whether a human or machine wrote the appeal.",
    Inches(2), Inches(4), Inches(9), Inches(0.8), font_size=20, color=GRAY)

# ============== SLIDE 17: Thank You ==============
slide = add_dark_slide(prs)

title = add_title_text(slide, "Thank You", Inches(0.8), font_size=56, color=LIGHT)
title.left = Inches(0)
title.width = prs.slide_width
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title.top = Inches(2)

sub = add_body_text(slide, "Questions?", Inches(0), Inches(3), prs.slide_width, Inches(0.8), font_size=32, color=ACCENT)
sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

info = add_body_text(slide, "Salah-din Mrait\ns.mrait@tilburguniversity.edu\n\nSupervisor: John Caffier",
    Inches(0), Inches(4.5), prs.slide_width, Inches(2), font_size=18, color=GRAY)
info.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('thesis_presentation_v2.pptx')
print("PowerPoint presentation created: thesis_presentation_v2.pptx")
